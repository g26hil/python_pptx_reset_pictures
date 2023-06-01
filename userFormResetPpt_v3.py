#pop a userform and open a ppt to reset the pictures
#

import tkinter as tk
from tkinter import ttk
from tkinter import filedialog

import time
from math import sqrt

import pptx
from pptx import Presentation
from pptx import slide
from pptx.util import Inches,Cm,Emu


########################################################################
#pretend we're const
MAXROWS = 3
SIZEFIX = 914400
ORI2CM = 360000
RNG_TOP = 2.09
RGN_LEFT = 0.7
RNG_HEIGHT = 16.43
RNG_WIDTH = 24.09
RNG_GAP = 0.1
RNG_USED = 0.6
LISTJOIN = ';'
PIC_NUM_MIN = 2
NO_RESET_PIC = 'NO RESET PICTURE'

##########################################################################
#define a class, get a file name, open a ppt, deal with all the pictures
class Reset_Prst:
    def __init__(self, fileName):
        self.fileName = fileName
        try:
            self.prs = Presentation(fileName)
        except:
            print("Error: File not found")
            return        
        #extremely important
        self.logName = self.new_log(self.fileName)
        #temporally
        self.RNG_WIDTH = RNG_WIDTH
        self.RNG_HEIGHT = RNG_HEIGHT
        self.rngProp = RNG_WIDTH/RNG_HEIGHT
        #refresh every slide
        self.clear_all()
        #default switch
        self.dealOver = False
        self.horiMid = False
        self.limitRow = True
        self.noRest= False

    def reset_pics(self):
        if self.rowNum > 1:
            return self.pic_to_multi()
        else:
            rtNum = self.log_str('just 1 row :',self.logName)
            self.pic_to1()
            return self.RNG_WIDTH
        
    def pic_to1(self):
        #present width of the first row
        lenOfRow = 0
        #the sum of all pic's  proportion
        picPropSum = sum(self.picsProp)
        #thin enough, just use the whole height, or use the whole width
        if picPropSum > self.rngProp:
            stdHeight = self.RNG_WIDTH / picPropSum
        else:
            stdHeight = self.RNG_HEIGHT
        #reset each pic
        for i in range(len(self.pics)):
            pic = self.pics[i]
            picLeft = RGN_LEFT + lenOfRow
            picTop = RNG_TOP
            picWidth = stdHeight * self.picsProp[i]
            lenOfRow += (picWidth + RNG_GAP)
            # deal with 90 and 270
            self.reset_1pic(pic,stdHeight,picWidth,picLeft,picTop,pic.rotation)

    def pic_to_multi(self):
        for i in range(self.rowNum):
            self.lenOfRow.append(0)
        dbStdHeight = self.RNG_HEIGHT/self.rowNum*2 + RNG_GAP
        stdHeight = self.RNG_HEIGHT/self.rowNum
        rowPst = 0
        #double rows
        while(len(self.dbPics) > 0):
            pic = self.dbPics.pop(0)
            picWidth = dbStdHeight * self.dbPicsProp.pop(0)
            picLeft = RGN_LEFT + self.lenOfRow[0]
            picTop =  RNG_TOP

            self.lenOfRow[0] += (picWidth + RNG_GAP)
            self.lenOfRow[1] += (picWidth + RNG_GAP)

            rtNum = self.log_str(self.dbPicsName.pop(0),self.logName)
            # deal with 90 and 270
            self.reset_1pic(pic,dbStdHeight,picWidth,picLeft,picTop,self.dbPicsRoation.pop(0))

        #single row
        while(len(self.pics) > 0):
            pic = self.pics.pop(0)
            picWidth = stdHeight * self.picsProp.pop(0)

            rtNum = self.log_str(self.picsName.pop(0),self.logName)

            rowPst = self.lenOfRow.index(min(self.lenOfRow))
            picLeft = RGN_LEFT + self.lenOfRow[rowPst]
            picTop =  RNG_TOP + rowPst * (stdHeight + RNG_GAP)

            self.lenOfRow[rowPst] += (picWidth + RNG_GAP)
            # deal with 90 and 270
            self.reset_1pic(pic,stdHeight,picWidth,picLeft,picTop,self.picsRoation.pop(0))

        return max(self.lenOfRow)

    def clear_all(self):
        #original pics info or which one will be put in one row finally
        self.pics=[]
        self.picsName=[]
        self.picsProp=[]
        self.picsRoation=[]
        #pics which will be put in multi row
        self.dbPics=[]
        self.dbPicsName=[] 
        self.dbPicsProp=[]
        self.dbPicsRoation=[]
        #the number of reset rows
        self.rowNum = 1
        #the width of each row, always change when pics be reseted      
        self.lenOfRow = []


    def visitPrt(self):
        for i in range(len(self.prs.slides)):
            self.sld = self.prs.slides[i]
            rtNum = self.log_str(str(i),self.logName)
            #
            self.visitSld()

    def visitSld(self):
        self.clear_all()
        for i in range(len(self.sld.shapes)):
            self.shp = self.sld.shapes[i]
            #if find a paticular string, just return -1
            if self.noRest:
                if self.shp.has_text_frame:
                    if self.shp.text == NO_RESET_PIC:
                        self.log_str(NO_RESET_PIC,self.logName)
                        return -1
            #collect all of pictures
            if self.shp.shape_type == 13:
                self.pics.append(self.shp)
                self.picsName.append(self.shp.name)
                self.picsRoation.append(self.shp.rotation)
                #deal with 90 and 270
                if self.shp.rotation == 90 or self.shp.rotation == 270:
                    self.picsProp.append(self.shp.height/self.shp.width)
                else:
                    self.picsProp.append(self.shp.width/self.shp.height)


        #make sure the number of piccures is not zero
        if len(self.pics)>0:
            rtNum = self.log_str('has pictures : ' + str(len(self.pics)),self.logName)
            #if the number of pictures is not too small, add rows, or just 1 row
            if len(self.pics) > PIC_NUM_MIN:
                self.rowNum = self.getRowNum()
                if self.limitRow:
                    self.rowNum = min(MAXROWS,self.rowNum)
            else:
                self.rowNum = 1
            rtNum = self.log_str('putin rowNum : ' + str(self.rowNum),self.logName)
            #find a way expand the picture, use space as posiblle as could
            if self.rowNum > 1:
                self.arrangePicToRow()
            #call subfuction to reset the pictures
            rtNum = self.reset_pics()
            #
            if rtNum > self.RNG_WIDTH and self.dealOver:
                scaleNum = rtNum/self.RNG_WIDTH
                self.log_str('over size : ' + str(rtNum) + ' and scale to small : ' + str(scaleNum),self.logName)
                self.deal_over(scaleNum)
            if rtNum < self.RNG_WIDTH and self.horiMid:
                pass


    def getRowNum(self):
        rowNums = 1
        picPropSum = sum(self.picsProp)
        while((self.rngProp * (rowNums**2)) < picPropSum):
            rowNums += 1
        return rowNums

    def arrangePicToRow(self):
        picPropSum = sum(self.picsProp)
        rngPropSum = self.rngProp * (self.rowNum**2)
        while((picPropSum/rngPropSum) < RNG_USED):
            #search the smallest proportion
            listPst = self.picsProp.index(min(self.picsProp))
            #calculate the proportion sum
            picPropSum += self.picsProp[listPst] * 3
            #prevent oversize
            if picPropSum  > rngPropSum:
                picPropSum -= self.picsProp[listPst] * 3
                break
            #move it to double list
            self.dbPicsName.append(self.picsName.pop(listPst))
            self.dbPicsProp.append(self.picsProp.pop(listPst))
            self.dbPics.append(self.pics.pop(listPst))
            self.dbPicsRoation.append(self.picsRoation.pop(listPst))

        rtNum = self.log_str('pictures : ' + str(picPropSum) + LISTJOIN + 'ranges : ' + str(rngPropSum),self.logName)
        rtNum = self.log_str('single line : ' , self.logName)
        rtNum = self.log_list(self.picsName,self.logName)
        rtNum = self.log_str('double line : ' , self.logName)
        rtNum = self.log_list(self.dbPicsName,self.logName)
        rtNum = self.log_str('end of arrangement' , self.logName)

    #create a txt named by the input file name and the real time, write a line of the input file name into it, return log name
    def new_log(self,fileName):
        #logName is a empty string
        logName=''
        fullPathList = fileName.split('.')
        if len(fullPathList) > 2:
            #group fullPathList into a string expcept the last item
            for i in range(0,len(fullPathList)-1):
                logName = logName + '.' +  fullPathList[i]
            #delelte the firs character of logNam
            logName = logName[1:]
        else:
            logName = fullPathList[0]
        #make a postfix for logName
        postfix = time.strftime('%Y%m%d_%H%M%S',time.localtime(time.time()))
        #make a log name based on fileName and real time
        logName = logName + '_' + postfix + '.txt'
        #create a txt file,and write some words in it
        with open(logName,'w') as f:
            f.write(fileName + '\n')
        #
        return logName


    #define a function to write a list into gived txt file ,which is seperated by ';'
    def log_list(self,list,logName):
        with open(logName,'a') as f:
            for item in list:
                f.write(str(item) + LISTJOIN)
            f.write('\n')
        return 0

    #define a fucntion to write a str into gived txt file
    def log_str(self,str,logName):
        with open(logName,'a') as f:
            f.write(str + '\n')
        return 0
    
    #define functions to deal 'with'
    def __exit__(self, exc_type, exc_value, traceback):
        self.prs.save(self.fileName)
        print("File reset")
    
    def __enter__(self):
        return self
    #define a function to get width and height of a slide of the gived prsentation,and return a list,and the size is based on inch
    def get_slide_size(self,prst):
        size = []
        size.append(Inches(prst.slide_width/Emu(1)).cm/SIZEFIX)
        size.append(Inches(prst.slide_height/Emu(1)).cm/SIZEFIX)
        return size
    #get the slide size，and reset the area for pictures   
    def resize_area(self):
        self.layout = self.get_slide_size(self.prs)
        self.log_str('layout:',self.logName)
        self.log_list(self.layout,self.logName)
        self.RNG_WIDTH = self.layout[0] - RGN_LEFT * 2
        self.RNG_HEIGHT = self.layout[1] - RNG_TOP - RGN_LEFT
        self.log_str('RNG_WIDTH:' + str(self.RNG_WIDTH) + LISTJOIN +' RNG_HEIGHT:' + str(self.RNG_HEIGHT),self.logName)
        self.rngProp = self.RNG_WIDTH / self.RNG_HEIGHT
        self.log_str('rngProp:' + str(self.rngProp),self.logName)

    def deal_over(self,scales):
        for i in range(len(self.sld.shapes)):
            self.shp = self.sld.shapes[i]
            #reset all of pictures
            if self.shp.shape_type == 13:
                self.shp.width = int(self.shp.width / scales)
                self.shp.height = int(self.shp.height / scales)
                self.shp.left = int((self.shp.left-Cm(RGN_LEFT)) / scales + Cm(RGN_LEFT))
                self.shp.top = int((self.shp.top-Cm(RNG_TOP)) / scales + Cm(RNG_TOP))

    def deal_revolved(self,pic,stdHeight,picWidth,picLeft,picTop):
        pic.width = Cm(stdHeight)
        pic.height = Cm(picWidth)
        revolvedX=picLeft+(picWidth/2)
        revolvedY=picTop+(stdHeight/2)
        pic.left = Cm(revolvedX-(stdHeight/2))
        pic.top = Cm(revolvedY-(picWidth/2))
    
    def deal_noRevolved(self,pic,stdHeight,picWidth,picLeft,picTop):
        pic.width = Cm(picWidth)
        pic.height = Cm(stdHeight)
        pic.left = Cm(picLeft)
        pic.top = Cm(picTop)

    def reset_1pic(self,pic,stdHeight,picWidth,picLeft,picTop,picRotation):
        if picRotation==90 or picRotation==270:
            self.log_str(pic.name + '--picRotation:' + str(picRotation),self.logName)
            self.deal_revolved(pic,stdHeight,picWidth,picLeft,picTop)
        else:
            self.deal_noRevolved(pic,stdHeight,picWidth,picLeft,picTop)


####################################################################################################################



class UserForm(tk.Frame):
    def __init__(self, parent, *args, **kwargs):
        tk.Frame.__init__(self, parent, *args, **kwargs)
        self.parent = parent

        if len(args)>0:
            self.arg0 = args[0]
            self.output = True

        self.user_form = tk.Frame(self.parent)
        self.user_form.pack()
        
        #caption of the frame
        self.user_form_label = tk.Label(self.user_form, text="Rest pictures of PPT")
        self.user_form_label.grid(row=0, column=0, columnspan=2, pady=10, padx=10, sticky="w")

        #a button to get a file
        self.user_form_button_get_file= tk.Button(self.user_form, text="选择一个PPT", command=self.get_file)
        self.user_form_button_get_file.grid(row=1, column=0, pady=10, padx=10, sticky="w")

        #textbox to show the file path
        self.user_form_textbox = tk.Text(self.user_form, height=5, width=60)
        self.user_form_textbox.grid(row=1, column=1, columnspan=2, pady=10, padx=10, sticky="w")

        #check box 1 to deal if pics out of range
        self.user_form_toggle_value = tk.BooleanVar(False)
        self.user_form_toggle_button = tk.Checkbutton(self.user_form, text="照片不能超出边界", variable=self.user_form_toggle_value)
        self.user_form_toggle_button.grid(row=2, column=0, pady=10, padx=10, sticky="w")

        #check box 2 to limits rows
        self.user_form_checkbox_value = tk.BooleanVar(False)
        self.user_form_checkbox = tk.Checkbutton(self.user_form, text="最多3行",variable=self.user_form_checkbox_value)
        self.user_form_checkbox.grid(row=3, column=0, pady=10, padx=10, sticky="w")
        
        #apply button, to save the userform, and destroy the frame
        self.user_form_button = tk.Button(self.user_form, text="提交执行",command=self.save_user_form)
        self.user_form_button.grid(row=3, column=1, pady=10, padx=10, sticky="w")

        #cancel button, to destroy the frame
        self.user_form_cancel_button = tk.Button(self.user_form, text="Cancel",command=self.parent.destroy)
        self.user_form_cancel_button.grid(row=3, column=2, pady=10, padx=10, sticky="w")

        #checkbox 3 to igone paticular slide
        self.user_form_checkbox_stop_value = tk.BooleanVar(False)
        self.user_form_checkbox_stop = tk.Checkbutton(self.user_form, text="忽略整页(NO RESET PICTURE)",variable=self.user_form_checkbox_stop_value)
        self.user_form_checkbox_stop.grid(row=4, column=0, pady=10, padx=10, sticky="w")

    def save_user_form(self):
        if self.output:
            self.arg0.append(self.user_form_textbox.get("1.0", "end-1c"))   #file name
            self.arg0.append(self.user_form_toggle_value.get())             #deal over range
            self.arg0.append(self.user_form_checkbox_value.get())           #limit rows
            self.arg0.append(self.user_form_checkbox_stop_value.get())     #ignor slide
        else:
            print("User Form Saved")
            print(self.user_form_textbox.get("1.0", "end-1c"))
            print(self.user_form_checkbox_value.get())
            print(self.user_form_toggle_value.get())
            print(self.user_form_checkbox_stop_value.get())
        self.parent.destroy()

    #get a file name and put in the textbox
    def get_file(self):
        strOfFile = self.pop_file_dialog("Select a file to open", "C:\\Users\\Jacob\\Documents\\GitHub\\Python-Projects\\tkinter-projects\\file-dialog-test", [("PPT Files", "*.pptx"), ("All Files", "*.*")])
        self.user_form_textbox.delete("1.0", "end")
        self.user_form_textbox.insert("1.0", strOfFile)

    # pop out a system msgbox to pich up a file
    def pop_file_dialog(self,title, directory, file_type):
        title=title,
        file_path = filedialog.askopenfilename(
            initialdir=directory,
            filetypes=file_type
        )
        return file_path


#################################################################################################
#   main pro
if __name__== "__main__":
    paraList = []
    u1=UserForm(tk.Tk(),paraList)
    u1.mainloop()
    #after the userform is gone
    print("it is all over!")
    if len(paraList)>0:
        for i in paraList:
            print(i)
        with Reset_Prst(paraList[0]) as p1:
            p1.resize_area()
            p1.dealOver=paraList[1]
            p1.limitRow=paraList[2]
            p1.noRest=paraList[3]
            p1.visitPrt()